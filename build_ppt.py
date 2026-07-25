# -*- coding: utf-8 -*-
"""生成「捕星少年」PPT：含一键启动按钮（点击运行启动捕星少年.bat）"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

CREAM  = RGBColor(0xFF, 0xFA, 0xF0)
CORAL  = RGBColor(0xFF, 0x8C, 0x7A)
SKY    = RGBColor(0x6B, 0xB6, 0xE0)
YELLOW = RGBColor(0xFF, 0xC8, 0x57)
INK    = RGBColor(0x2D, 0x2A, 0x26)
SOFT   = RGBColor(0x5C, 0x51, 0x47)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

APP_DIR = r"C:\Users\xiaos\WorkBuddy\2026-07-24-20-56-58\app"
BAT = os.path.join(APP_DIR, "启动捕星少年.bat")
OUT = os.path.join(APP_DIR, "捕星少年.pptx")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def textbox(slide, l, t, w, h, lines, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, item in enumerate(lines):
        txt, size, color = item[0], item[1], item[2]
        bold = item[3] if len(item) > 3 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
        r.font.name = "Microsoft YaHei"
    return tb

def button(slide, l, t, w, h, label, link):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = CORAL
    shp.line.color.rgb = CORAL
    tf = shp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Microsoft YaHei"
    # 点击打开启动脚本
    shp.click_action.hyperlink.address = link
    return shp

# ---------- Slide 1: 封面 + 启动按钮 ----------
s = prs.slides.add_slide(BLANK); bg(s, CREAM)
textbox(s, 1.2, 1.6, 10.9, 1.6, [("捕星少年", 60, CORAL, True)])
textbox(s, 1.2, 3.1, 10.9, 0.9, [("暑假作业打卡 · 家庭版", 26, SKY, True)])
button(s, 4.4, 4.4, 4.5, 1.2, "▶ 打开捕星少年", BAT)
textbox(s, 1.2, 5.9, 10.9, 0.8,
        [("点击按钮会自动启动服务并在浏览器打开应用（渲染更清晰）", 14, SOFT, False)])

# ---------- Slide 2: 怎么用 ----------
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
textbox(s, 0.8, 0.5, 11.7, 1.0, [("怎么用", 36, CORAL, True)])
steps = [
    ("1. 在封面页点击「▶ 打开捕星少年」按钮", 20, INK),
    ("2. 等待浏览器自动弹出捕星少年（首次会后台启动服务）", 20, INK),
    ("3. 用账号登录：家长 Mum（手动输入）/ 孩子 Damon、Lemon", 20, INK),
    ("4. 孩子端录入作业、每天打卡攒星星；家长端看进度、设奖励", 20, INK),
    ("5. 用完直接关浏览器即可，服务在后台继续运行", 20, INK),
]
textbox(s, 1.0, 1.8, 11.3, 4.8, steps, align=PP_ALIGN.LEFT)

# ---------- Slide 3: 账号一览 ----------
s = prs.slides.add_slide(BLANK); bg(s, CREAM)
textbox(s, 0.8, 0.5, 11.7, 1.0, [("账号一览", 36, SKY, True)])
acc = [
    ("家长端  Mum", 22, CORAL),
    ("已隐藏，不在登录页显示；登录时手动输入用户名 Mum，密码 20260725", 16, SOFT),
    ("孩子端  Damon", 22, INK),
    ("用户名 Damon / 密码 2013", 16, SOFT),
    ("孩子端  Lemon", 22, INK),
    ("用户名 Lemon / 密码 2016", 16, SOFT),
]
textbox(s, 1.0, 1.8, 11.3, 4.8, acc, align=PP_ALIGN.LEFT)

# ---------- Slide 4: 注意事项 ----------
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
textbox(s, 0.8, 0.5, 11.7, 1.0, [("注意事项", 36, CORAL, True)])
notes = [
    ("· 启动按钮依赖本机已安装 Node.js（已具备），并会自动在后台运行服务", 18, INK),
    ("· 数据全部存在本机 app/data.json，不会上传到任何网络，隐私安全", 18, INK),
    ("· 服务在后台运行时，关掉浏览器不会停止；重启电脑后需再点一次按钮", 18, INK),
    ("· 想要「在幻灯片里直接显示网页」可试 PowerPoint 的 Web Viewer 插件，", 18, INK),
    ("  但旧版网页引擎可能排版不全，故推荐用浏览器方式使用", 18, INK),
]
textbox(s, 1.0, 1.8, 11.3, 4.8, notes, align=PP_ALIGN.LEFT)

prs.save(OUT)
print("saved:", OUT)
